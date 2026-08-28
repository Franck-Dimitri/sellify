import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG_DARK = RGBColor(15, 23, 42)    # #0F172A (Dark Slate)
    COLOR_BG_LIGHT = RGBColor(248, 250, 252) # #F8FAFC (Light Gray)
    COLOR_PRIMARY = RGBColor(234, 179, 8)   # #EAB308 (Brand Yellow)
    COLOR_SECONDARY = RGBColor(16, 185, 129) # #10B981 (Emerald)
    COLOR_TEXT_DARK = RGBColor(15, 23, 42)  # #0F172A
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_CARD_BG = RGBColor(255, 255, 255)
    COLOR_CARD_BORDER = RGBColor(226, 232, 240)

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category="SELLIFY.ME STRATÉGIE"):
        # Top Accent Line
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_PRIMARY
        shape.line.fill.background()

        # Category
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.3))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

        # Title
        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.6))
        tf2 = tx_box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_DARK

    def set_slide_background(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    # =========================================================================
    # SLIDE 1 : COVER SLIDE (DARK THEME)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, COLOR_BG_DARK)

    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(30, 41, 59)
    badge.line.color.rgb = COLOR_PRIMARY
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ DOSSIER DE PRÉSENTATION OFFICIEL DU PROJET"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    tx = s1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "SELLIFY.ME"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf.add_paragraph()
    p2.text = "La Plateforme E-Commerce de Confiance pour l'Afrique"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_PRIMARY

    tx2 = s1.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.8))
    tf2 = tx2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Séquestre Escrow Mobile Money · Logistique IA & Télémétrie · SaaS Marchand & Smart-Links · Moteur IA 1.2 Flash"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(203, 213, 225)

    card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card.line.color.rgb = RGBColor(51, 65, 85)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = "Stack : Laravel 11 · React 19 · Inertia.js · PostgreSQL 16 · OSRM · Sellify AI 1.2 Flash"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Périmètre de déploiement : Phase 1 : Cameroun (Douala & Yaoundé) | Phase 2 : Zone CEMAC & CEDEAO"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(148, 163, 184)

    # =========================================================================
    # SLIDE 2 : EXECUTIVE SUMMARY & VISION
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, COLOR_BG_LIGHT)
    add_header(s2, "Executive Summary & Vision du Projet", "VISION STRATÉGIQUE")

    cards_data = [
        ("1. Problème Fondamental", "Le e-commerce africain est paralysé par le manque de confiance, l'échec du Cash on Delivery (>35% de refus) et l'absence d'adresses postales formelles.", COLOR_PRIMARY),
        ("2. Solution Sellify.me", "Une plateforme intégrée combinant le séquestre Escrow Mobile Money (OM/MoMo), un guidage GPS OSRM en temps réel et des Smart-Links de vente sociale.", COLOR_SECONDARY),
        ("3. Impact Économique", "Professionnalisation de centaines de milliers de micro-marchands informels, réduction des litiges à 0% et sécurité totale des encaissements.", RGBColor(59, 130, 246)),
    ]

    for idx, (ctitle, cdesc, ccolor) in enumerate(cards_data):
        cx = Inches(0.8 + idx * 4.0)
        cshape = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.6), Inches(3.7), Inches(5.0))
        cshape.fill.solid()
        cshape.fill.fore_color.rgb = COLOR_CARD_BG
        cshape.line.color.rgb = COLOR_CARD_BORDER

        bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, Inches(1.6), Inches(3.7), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ccolor
        bar.line.fill.background()

        tx = s2.shapes.add_textbox(cx + Inches(0.2), Inches(1.9), Inches(3.3), Inches(4.5))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = cdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 3 : LES 4 PILIERS DE L'ÉCOSYSTÈME
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, COLOR_BG_LIGHT)
    add_header(s3, "Les 4 Piliers Fondateurs de l'Écosystème", "ÉCOSYSTÈME SELLIFY")

    pillars = [
        ("1. Confiance Totale (Escrow Core)", "Séquestre numérique inviolable sur Orange Money / MTN MoMo. Déblocage conditionné à la validation de l'acheteur.", COLOR_PRIMARY),
        ("2. Logistique IA & Navigation GPS", "Guidage GPS routier embarqué OSRM, télémétrie en direct et calcul continu d'ETA sans quitter l'application.", COLOR_SECONDARY),
        ("3. Suite SaaS Marchande & Smart-Links", "Gestion multi-boutiques, catalogue avancé et tunnels de commande 1-clic pour WhatsApp, Facebook et TikTok.", RGBColor(99, 102, 241)),
        ("4. Copilote Universel Sellify AI 1.2 Flash", "Assistant conversationnel multi-rôles avec reconnaissance & synthèse vocale, détection anti-fraude et aide au pricing.", RGBColor(14, 165, 233)),
    ]

    for idx, (ptitle, pdesc, pcolor) in enumerate(pillars):
        py = Inches(1.6 + idx * 1.25)
        pshape = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), py, Inches(11.7), Inches(1.1))
        pshape.fill.solid()
        pshape.fill.fore_color.rgb = COLOR_CARD_BG
        pshape.line.color.rgb = COLOR_CARD_BORDER

        stripe = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), py, Inches(0.15), Inches(1.1))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = pcolor
        stripe.line.fill.background()

        tx = s3.shapes.add_textbox(Inches(1.2), py + Inches(0.1), Inches(11.0), Inches(0.9))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ptitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = pdesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 4 : LE SÉQUESTRE ESCROW MOBILE MONEY
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, COLOR_BG_LIGHT)
    add_header(s4, "Le Moteur Transactionnel d'Escrow Mobile Money", "SÉCURITÉ FINANCIÈRE")

    steps_data = [
        ("Étape 1 : Commande", "Le client paie par Orange Money ou MTN MoMo. Les fonds sont bloqués sur le compte séquestre du vendeur (pending_balance).", "Statut: 'escrow_held'"),
        ("Étape 2 : Livraison", "Le vendeur prépare le colis et le confie au livreur. Le client suit le colis en temps réel sur la carte OSRM.", "Code Secret OTP généré"),
        ("Étape 3 : Déblocage", "Le livreur valide le code OTP secret 6 chiffres fourni par le client. Les fonds basculent sur le solde disponible du vendeur.", "Statut: 'released'"),
        ("Étape 4 : Retrait", "Le vendeur retire immédiatement ses gains vers son compte Mobile Money ou Banque sous 24h.", "Solde Disponible"),
    ]

    for idx, (stitle, sdesc, sbadge) in enumerate(steps_data):
        sx = Inches(0.8 + idx * 2.95)
        sshape = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, Inches(1.7), Inches(2.8), Inches(4.8))
        sshape.fill.solid()
        sshape.fill.fore_color.rgb = COLOR_CARD_BG
        sshape.line.color.rgb = COLOR_CARD_BORDER

        badge = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx + Inches(0.2), Inches(1.9), Inches(2.4), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(254, 240, 138)
        badge.line.color.rgb = COLOR_PRIMARY
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = sbadge
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = RGBColor(113, 63, 18)
        p.alignment = PP_ALIGN.CENTER

        tx = s4.shapes.add_textbox(sx + Inches(0.2), Inches(2.4), Inches(2.4), Inches(3.9))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 5 : LOGISTIQUE ROUTIÈRE & GPS TURN-BY-TURN
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, COLOR_BG_LIGHT)
    add_header(s5, "Logistique Urbaine & Guidage GPS Embarqué", "LOGISTIQUE GÉOLOCALISÉE")

    l_items = [
        ("Moteur OSRM Intégré", "Tracé routier suivant fidèlement les boulevards, avenues et carrefours réels de Douala et Yaoundé sans jamais sortir vers une application externe."),
        ("Double Polyligne Dynamique", "Visualisation claire du chemin parcouru en vert et du segment restant en jaune pointillé avec calcul continu de distance (km) et d'ETA (minutes)."),
        ("Double Clôture Sécurisée", "La livraison ne peut être clôturée qu'avec la saisie du code OTP secret client à 6 chiffres + la signature tactile sur écran + photo preuve."),
        ("Système de Dispatch IA", "Attribution optimale des courses selon la géolocalisation des chauffeurs, l'encombrement et la disponibilité."),
    ]

    for idx, (ltitle, ldesc) in enumerate(l_items):
        lx = Inches(0.8 + (idx % 2) * 6.0)
        ly = Inches(1.8 + (idx // 2) * 2.5)

        lshape = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ly, Inches(5.7), Inches(2.2))
        lshape.fill.solid()
        lshape.fill.fore_color.rgb = COLOR_CARD_BG
        lshape.line.color.rgb = COLOR_CARD_BORDER

        tx = s5.shapes.add_textbox(lx + Inches(0.3), ly + Inches(0.2), Inches(5.1), Inches(1.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"📍 {ltitle}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = ldesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 6 : SUITE MARCHAND & SMART-LINKS
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, COLOR_BG_LIGHT)
    add_header(s6, "Suite SaaS Vendeurs & Smart-Links de Vente Sociale", "COMMERCE MARCHAND")

    s_items = [
        ("Multi-Boutiques Centralisé", "Supervision de plusieurs enseignes commerciales depuis un tableau de bord unique avec isolation des stocks."),
        ("Smart-Links 1-Clic", "Génération de liens de paiement rapides à partager directement sur WhatsApp, TikTok et Instagram avec encaissement Mobile Money intégré."),
        ("Portefeuille & Grand Livre", "Graphiques interactifs d'évolution des flux, exports de transactions en CSV/Excel et impression de Relevés de Compte Officiels certifiés (PDF)."),
        ("Moteur Marketing & Remises", "Création de codes promo personnalisés, remises dégressives sur le volume d'achat et bannières publicitaires."),
    ]

    for idx, (stitle, sdesc) in enumerate(s_items):
        sx = Inches(0.8 + (idx % 2) * 6.0)
        sy = Inches(1.8 + (idx // 2) * 2.5)

        sshape = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, Inches(5.7), Inches(2.2))
        sshape.fill.solid()
        sshape.fill.fore_color.rgb = COLOR_CARD_BG
        sshape.line.color.rgb = COLOR_CARD_BORDER

        tx = s6.shapes.add_textbox(sx + Inches(0.3), sy + Inches(0.2), Inches(5.1), Inches(1.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"🛍️ {stitle}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 7 : GOUVERNANCE & ARBITRAGE SUPERADMIN
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, COLOR_BG_LIGHT)
    add_header(s7, "Gouvernance, Modération KYC & Arbitrage SuperAdmin", "CONTRÔLE & CONFORMITÉ")

    g_items = [
        ("Vérification KYC Multi-Niveaux", "Contrôle strict des pièces d'identité, registres de commerce (RCCM) et cartes grises des marchands et des livreurs partenaires."),
        ("Centre d'Arbitrage des Litiges", "Interface dédiée de médiation permettant la libération des fonds vers le vendeur ou le remboursement immédiat à l'acheteur en cas de contestation."),
        ("Supervision des Flux Financiers", "Tableau de bord de suivi en temps réel du volume total sous séquestre, du solde en attente et des commissions de plateforme."),
        ("Journal d'Audit Immuable", "Enregistrement cryptographique et horodaté de chaque action système dans la table activity_logs avec conservation des adresses IP."),
    ]

    for idx, (gtitle, gdesc) in enumerate(g_items):
        gx = Inches(0.8 + (idx % 2) * 6.0)
        gy = Inches(1.8 + (idx // 2) * 2.5)

        gshape = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gx, gy, Inches(5.7), Inches(2.2))
        gshape.fill.solid()
        gshape.fill.fore_color.rgb = COLOR_CARD_BG
        gshape.line.color.rgb = COLOR_CARD_BORDER

        tx = s7.shapes.add_textbox(gx + Inches(0.3), gy + Inches(0.2), Inches(5.1), Inches(1.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"🛡️ {gtitle}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = gdesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 8 : MOTEUR SELLIFY AI 1.2 FLASH
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, COLOR_BG_LIGHT)
    add_header(s8, "Moteur Sellify AI 1.2 Flash : Le Copilote Universel", "INTELLIGENCE ARTIFICIELLE")

    ai_roles = [
        ("👤 Copilote Acheteur", "Recherche vocale de produits, localisation instantanée de colis, vérification de l'OTP et assistance aux litiges."),
        ("🏬 Copilote Vendeur", "Optimisation des prix de vente, génération automatique de descriptions produits engageantes et simulation de marge."),
        ("🛵 Copilote Livreur", "Guidage routier vocal naturel, alertes embouteillages en temps réel et bilan des gains journaliers."),
        ("🛡️ Copilote SuperAdmin", "Surveillance anti-fraude, détection des anomalies sur les flux Escrow et aide à l'arbitrage des litiges."),
    ]

    for idx, (rtitle, rdesc) in enumerate(ai_roles):
        rx = Inches(0.8 + (idx % 2) * 6.0)
        ry = Inches(1.8 + (idx // 2) * 2.5)

        rshape = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, ry, Inches(5.7), Inches(2.2))
        rshape.fill.solid()
        rshape.fill.fore_color.rgb = COLOR_CARD_BG
        rshape.line.color.rgb = COLOR_CARD_BORDER

        tx = s8.shapes.add_textbox(rx + Inches(0.3), ry + Inches(0.2), Inches(5.1), Inches(1.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = rtitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = rdesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 9 : MODÈLE ÉCONOMIQUE & RENTABILITÉ
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, COLOR_BG_LIGHT)
    add_header(s9, "Modèle Économique & Flux de Monétisation", "BUSINESS MODEL")

    bm_data = [
        ("1. Commissions Escrow (3%)", "Prélèvement automatique de 3% sur chaque vente finalisée via la marketplace et les Smart-Links."),
        ("2. Abonnements SaaS Vendeur", "Formules Starter (Gratuit), Pro (15 000 FCFA/mois) et Enterprise (45 000 FCFA/mois)."),
        ("3. Marge Logistique (10-15%)", "Commission prélevée sur les frais de livraison acheminés par les chauffeurs partenaires."),
        ("4. Services Premium & Visibilité", "Options de mise en avant publicitaire de boutiques et badges de certification prioritaires."),
    ]

    for idx, (bmtitle, bmdesc) in enumerate(bm_data):
        bx = Inches(0.8 + (idx % 2) * 6.0)
        by = Inches(1.8 + (idx // 2) * 2.5)

        bshape = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(5.7), Inches(2.2))
        bshape.fill.solid()
        bshape.fill.fore_color.rgb = COLOR_CARD_BG
        bshape.line.color.rgb = COLOR_CARD_BORDER

        tx = s9.shapes.add_textbox(bx + Inches(0.3), by + Inches(0.2), Inches(5.1), Inches(1.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💰 {bmtitle}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = bmdesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 10 : ARCHITECTURE TECHNIQUE & SÉCURITÉ
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, COLOR_BG_LIGHT)
    add_header(s10, "Architecture Technique & Normes de Qualité", "INGÉNIERIE LOGICIELLE")

    tech_specs = [
        ("Backend Framework", "Laravel 11.x avec Service Layer découplé (EscrowService, LogisticsIAService)."),
        ("Frontend Réactif", "React 19 + Inertia.js (Single Page Application fluide sans latence de requêtes REST superflues)."),
        ("Base de Données", "PostgreSQL 16 avec modélisation relationnelle stricte, indexation géographique et intégrité référentielle."),
        ("Cartographie & Télémétrie", "Leaflet.js + OpenStreetMap + Serveur de calcul d'itinéraires OSRM haute performance."),
        ("Tests & Validation", "47/47 tests automatisés validés sous PHPUnit (255 assertions, 100% de taux de succès)."),
        ("Sécurité & Conformité", "Hachage Argon2id, tokens Sanctum, vérification KYC multi-niveaux et audit trail immuable."),
    ]

    for idx, (ts_title, ts_desc) in enumerate(tech_specs):
        tx_pos = Inches(0.8 + (idx % 3) * 4.0)
        ty_pos = Inches(1.8 + (idx // 3) * 2.5)

        tshape = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx_pos, ty_pos, Inches(3.7), Inches(2.2))
        tshape.fill.solid()
        tshape.fill.fore_color.rgb = COLOR_CARD_BG
        tshape.line.color.rgb = COLOR_CARD_BORDER

        tx = s10.shapes.add_textbox(tx_pos + Inches(0.2), ty_pos + Inches(0.15), Inches(3.3), Inches(1.9))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ts_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

        p2 = tf.add_paragraph()
        p2.text = ts_desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 11 : FEUILLE DE ROUTE STRATÉGIQUE (ROADMAP)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11, COLOR_BG_LIGHT)
    add_header(s11, "Feuille de Route & Déploiement Panafricain", "EXPANSION STRATÉGIQUE")

    phases = [
        ("Phase 1 : 2026 Q3 (Actuel)", "• Lancement officiel Douala & Yaoundé\n• Escrow Orange Money & MTN MoMo\n• Copilote Sellify AI 1.2 Flash\n• Smart-Links de vente sociale", COLOR_PRIMARY),
        ("Phase 2 : 2026 Q4", "• Déploiement Apps Mobiles Android & iOS\n• Extension dans les villes secondaires\n• Hubs de distribution express\n• Partenariats logistiques régionaux", COLOR_SECONDARY),
        ("Phase 3 : 2027 Q1 - Q2", "• Expansion zone CEMAC (Gabon, Congo, Tchad)\n• Hubs de stockage urbain partagé\n• Passerelle de paiement transfrontalière\n• Extension vers la zone CEDEAO", RGBColor(59, 130, 246)),
    ]

    for idx, (phtitle, phdesc, phcolor) in enumerate(phases):
        phx = Inches(0.8 + idx * 4.0)
        phshape = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phx, Inches(1.8), Inches(3.7), Inches(4.7))
        phshape.fill.solid()
        phshape.fill.fore_color.rgb = COLOR_CARD_BG
        phshape.line.color.rgb = COLOR_CARD_BORDER

        bar = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, phx, Inches(1.8), Inches(3.7), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = phcolor
        bar.line.fill.background()

        tx = s11.shapes.add_textbox(phx + Inches(0.2), Inches(2.1), Inches(3.3), Inches(4.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phtitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = phdesc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 12 : CONCLUSION & CONTACT
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12, COLOR_BG_DARK)

    bar = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    tx = s12.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(3.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BÂTIR LA CONFIANCE DU COMMERCE AFRICAIN"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Sellify.me est prêt à transformer le potentiel commercial de l'Afrique en une réalité économique sécurisée et prospère."
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_PRIMARY
    p2.alignment = PP_ALIGN.CENTER

    card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(4.6), Inches(7.333), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card.line.color.rgb = RGBColor(51, 65, 85)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = "CONTACT DIRECTION & RELATIONS INVESTISSEURS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Email : direction@sellify.me · contact@sellify.me\nSiège Social : Boulevard de la Liberté, Akwa, Douala — Cameroun\nSite Web : https://sellify.me"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(226, 232, 240)
    p2.alignment = PP_ALIGN.CENTER

    output_path = "/home/mr-dims-tech/developpement/developpement_laravel/mr_dims/sellify/Presentation_Projet_Sellify.pptx"
    prs.save(output_path)
    print(f"PowerPoint Presentation successfully updated at: {output_path}")

if __name__ == "__main__":
    create_presentation()
